"""
Document Processor - Handles processing of various document types.
Supports PDF, DOCX, PPTX, images, audio, and video.
"""
import asyncio
import mimetypes
import hashlib
from pathlib import Path
from typing import List, Optional, Dict, Any, Tuple
from dataclasses import dataclass, field
import uuid
import shutil

from app.core.config import settings
from app.services.embedding_manager import EmbeddingManager
from app.services.vector_store import VectorStoreManager, Document


@dataclass
class ProcessedDocument:
    """Processed document with extracted content and metadata."""
    id: str
    filename: str
    file_type: str
    content: str
    chunks: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    page_count: int = 0
    word_count: int = 0


class DocumentProcessor:
    """
    Processes various document types and stores them in the vector store.
    Supports: PDF, DOCX, PPTX, TXT, MD, images (OCR), audio, video.
    """

    def __init__(
        self,
        embedding_manager: EmbeddingManager,
        vector_store: VectorStoreManager,
    ):
        self.embedding_manager = embedding_manager
        self.vector_store = vector_store
        self.documents_dir = settings.documents_dir

    async def process_file(
        self,
        file_path: Path,
        original_filename: str,
    ) -> ProcessedDocument:
        """
        Process a file and store in vector store.
        Returns processed document with metadata.
        """
        file_type = self._get_file_type(file_path, original_filename)

        # Extract content based on file type
        if file_type == "pdf":
            content, metadata = await self._process_pdf(file_path)
        elif file_type == "docx":
            content, metadata = await self._process_docx(file_path)
        elif file_type == "pptx":
            content, metadata = await self._process_pptx(file_path)
        elif file_type in ["txt", "md", "json", "csv"]:
            content, metadata = await self._process_text(file_path)
        elif file_type in ["png", "jpg", "jpeg", "gif", "webp"]:
            content, metadata = await self._process_image(file_path)
        elif file_type in ["mp3", "wav", "m4a", "flac", "ogg"]:
            content, metadata = await self._process_audio(file_path)
        elif file_type in ["mp4", "avi", "mov", "mkv", "webm"]:
            content, metadata = await self._process_video(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        # Generate document ID
        doc_id = str(uuid.uuid4())

        # Calculate word count
        word_count = len(content.split())

        # Create chunks for vector store
        chunks = self._create_chunks(content)

        # Store original file
        stored_path = await self._store_file(file_path, original_filename, doc_id)

        # Create processed document
        processed = ProcessedDocument(
            id=doc_id,
            filename=original_filename,
            file_type=file_type,
            content=content,
            chunks=chunks,
            metadata={
                **metadata,
                "stored_path": str(stored_path),
                "original_filename": original_filename,
            },
            page_count=metadata.get("page_count", 1),
            word_count=word_count,
        )

        # Store chunks in vector store
        await self._store_chunks(processed)

        return processed

    def _get_file_type(self, file_path: Path, filename: str) -> str:
        """Determine file type from extension or MIME type."""
        extension = Path(filename).suffix.lower().lstrip(".")

        if extension:
            return extension

        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type:
            return mime_type.split("/")[-1]

        return "unknown"

    async def _process_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF using PyMuPDF."""
        import fitz  # PyMuPDF

        loop = asyncio.get_event_loop()

        def extract():
            doc = fitz.open(str(file_path))
            text_parts = []
            page_count = len(doc)

            for page_num in range(page_count):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{text}")

            doc.close()
            return "\n\n".join(text_parts), {"page_count": page_count}

        return await loop.run_in_executor(None, extract)

    async def _process_docx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX."""
        from docx import Document as DocxDocument

        loop = asyncio.get_event_loop()

        def extract():
            doc = DocxDocument(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs), {"paragraph_count": len(paragraphs)}

        return await loop.run_in_executor(None, extract)

    async def _process_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX."""
        from pptx import Presentation

        loop = asyncio.get_event_loop()

        def extract():
            prs = Presentation(str(file_path))
            text_parts = []
            slide_count = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

            return "\n\n".join(text_parts), {"slide_count": slide_count}

        return await loop.run_in_executor(None, extract)

    async def _process_text(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Process text files."""
        content = file_path.read_text(encoding="utf-8", errors="ignore")
        line_count = content.count("\n") + 1
        return content, {"line_count": line_count}

    async def _process_image(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract text from image using OCR."""
        try:
            import easyocr

            loop = asyncio.get_event_loop()

            def extract():
                reader = easyocr.Reader(["en"], gpu=False)
                results = reader.readtext(str(file_path))
                text = "\n".join([result[1] for result in results])
                return text, {"ocr_regions": len(results)}

            return await loop.run_in_executor(None, extract)

        except ImportError:
            # Fallback to pytesseract
            try:
                import pytesseract
                from PIL import Image

                loop = asyncio.get_event_loop()

                def extract():
                    image = Image.open(file_path)
                    text = pytesseract.image_to_string(image)
                    return text, {"ocr_method": "tesseract"}

                return await loop.run_in_executor(None, extract)

            except ImportError:
                return "[Image processing not available]", {"error": "No OCR library installed"}

    async def _process_audio(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Transcribe audio using Whisper."""
        try:
            from faster_whisper import WhisperModel

            loop = asyncio.get_event_loop()

            def transcribe():
                model = WhisperModel("base", device="cpu", compute_type="int8")
                segments, info = model.transcribe(str(file_path))

                text_parts = []
                for segment in segments:
                    text_parts.append(segment.text)

                return " ".join(text_parts), {
                    "duration": info.duration,
                    "language": info.language,
                }

            return await loop.run_in_executor(None, transcribe)

        except ImportError:
            return "[Audio transcription not available]", {"error": "Whisper not installed"}

    async def _process_video(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract audio from video and transcribe."""
        try:
            from moviepy.editor import VideoFileClip

            # Extract audio to temp file
            temp_audio = settings.cache_dir / f"temp_audio_{uuid.uuid4()}.wav"

            loop = asyncio.get_event_loop()

            def extract_audio():
                video = VideoFileClip(str(file_path))
                video.audio.write_audiofile(str(temp_audio), verbose=False, logger=None)
                duration = video.duration
                video.close()
                return duration

            duration = await loop.run_in_executor(None, extract_audio)

            # Transcribe audio
            text, audio_metadata = await self._process_audio(temp_audio)

            # Clean up temp file
            temp_audio.unlink(missing_ok=True)

            return text, {
                **audio_metadata,
                "video_duration": duration,
                "source": "video_audio_track",
            }

        except ImportError:
            return "[Video processing not available]", {"error": "MoviePy not installed"}

    def _create_chunks(
        self,
        content: str,
        chunk_size: int = None,
        overlap: int = None,
    ) -> List[str]:
        """Split content into overlapping chunks."""
        chunk_size = chunk_size or settings.chunk_size
        overlap = overlap or settings.chunk_overlap

        if len(content) <= chunk_size:
            return [content]

        chunks = []
        start = 0

        while start < len(content):
            end = start + chunk_size

            # Try to end at sentence boundary
            if end < len(content):
                # Look for sentence ending
                for sep in [". ", ".\n", "!\n", "?\n", "\n\n"]:
                    last_sep = content[start:end].rfind(sep)
                    if last_sep > chunk_size // 2:
                        end = start + last_sep + len(sep)
                        break

            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap

        return chunks

    async def _store_chunks(self, processed: ProcessedDocument):
        """Store document chunks in vector store."""
        documents = []

        for i, chunk in enumerate(processed.chunks):
            doc = Document(
                id=f"{processed.id}_chunk_{i}",
                content=chunk,
                metadata={
                    "document_id": processed.id,
                    "chunk_index": i,
                    "total_chunks": len(processed.chunks),
                    "filename": processed.filename,
                    "file_type": processed.file_type,
                },
            )
            documents.append(doc)

        await self.vector_store.add_documents(documents)

    async def _store_file(
        self,
        file_path: Path,
        original_filename: str,
        doc_id: str,
    ) -> Path:
        """Store original file in documents directory."""
        # Create subdirectory for document
        doc_dir = self.documents_dir / doc_id
        doc_dir.mkdir(parents=True, exist_ok=True)

        # Copy file
        dest_path = doc_dir / original_filename
        shutil.copy2(file_path, dest_path)

        return dest_path

    async def get_document(self, doc_id: str) -> Optional[ProcessedDocument]:
        """Retrieve a processed document by ID."""
        # Get first chunk to get metadata
        chunk = await self.vector_store.get_document(f"{doc_id}_chunk_0")
        if not chunk:
            return None

        # Get all chunks
        chunks = []
        chunk_index = 0
        while True:
            chunk_doc = await self.vector_store.get_document(f"{doc_id}_chunk_{chunk_index}")
            if not chunk_doc:
                break
            chunks.append(chunk_doc.content)
            chunk_index += 1

        return ProcessedDocument(
            id=doc_id,
            filename=chunk.metadata.get("filename", ""),
            file_type=chunk.metadata.get("file_type", ""),
            content="\n\n".join(chunks),
            chunks=chunks,
            metadata=chunk.metadata,
        )

    async def delete_document(self, doc_id: str) -> bool:
        """Delete a document and its chunks."""
        # Delete from vector store
        deleted_count = await self.vector_store.delete_by_metadata({"document_id": doc_id})

        # Delete stored file
        doc_dir = self.documents_dir / doc_id
        if doc_dir.exists():
            shutil.rmtree(doc_dir)

        return deleted_count > 0

    async def list_documents(self) -> List[Dict[str, Any]]:
        """List all processed documents."""
        documents = []

        # Scan documents directory
        if self.documents_dir.exists():
            for doc_dir in self.documents_dir.iterdir():
                if doc_dir.is_dir():
                    doc_id = doc_dir.name
                    # Get metadata from vector store
                    chunk = await self.vector_store.get_document(f"{doc_id}_chunk_0")
                    if chunk:
                        documents.append({
                            "id": doc_id,
                            "filename": chunk.metadata.get("filename", ""),
                            "file_type": chunk.metadata.get("file_type", ""),
                            "chunk_count": chunk.metadata.get("total_chunks", 0),
                        })

        return documents